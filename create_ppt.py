from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── COLOUR PALETTE ───────────────────────────────────────────────
BG_DARK      = RGBColor(0x0A, 0x0A, 0x0F)   # #0a0a0f  deep dark
BG_CARD      = RGBColor(0x11, 0x11, 0x1A)   # #11111a  card bg
GOLD         = RGBColor(0xC9, 0xA8, 0x4C)   # #c9a84c  primary gold
GOLD_LIGHT   = RGBColor(0xE8, 0xC9, 0x7A)   # #e8c97a  light gold
WHITE        = RGBColor(0xF0, 0xED, 0xE8)   # #f0ede8  off-white
GREY         = RGBColor(0x9B, 0x95, 0x92)   # #9b9592  muted grey
ACCENT_BLUE  = RGBColor(0x1A, 0x1A, 0x2E)   # dark blue accent


def hex_rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


# ─── HELPERS ──────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=16, color=WHITE, spacing=1.15, font_name="Calibri"):
    """lines = list of (text, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for text, bold, clr in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = clr if clr else color
        run.font.name = font_name
    return txBox


def gold_bar(slide, top=0.55, height=0.06):
    """Thin gold accent bar near top"""
    add_rect(slide, 0, top, 13.33, height, fill_color=GOLD)


def slide_header(slide, title, subtitle=None):
    """Standard slide top area"""
    gold_bar(slide)
    add_text(slide, title,
             left=0.5, top=0.7, width=12, height=0.6,
             font_size=30, bold=True, color=GOLD_LIGHT,
             align=PP_ALIGN.LEFT, font_name="Calibri Light")
    if subtitle:
        add_text(slide, subtitle,
                 left=0.5, top=1.25, width=12, height=0.35,
                 font_size=14, bold=False, color=GREY,
                 align=PP_ALIGN.LEFT)


def add_bullet_card(slide, left, top, width, height, title, bullets,
                    title_size=15, bullet_size=13):
    """Card with title and bullet list"""
    add_rect(slide, left, top, width, height,
             fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1.2))
    add_text(slide, title,
             left=left+0.12, top=top+0.1, width=width-0.2, height=0.35,
             font_size=title_size, bold=True, color=GOLD)
    lines = []
    for b in bullets:
        lines.append(("  •  " + b, False, WHITE))
    add_multiline(slide, lines,
                  left=left+0.12, top=top+0.45, width=width-0.2,
                  height=height-0.55, font_size=bullet_size, color=WHITE)


def logo_text(slide):
    add_text(slide, "Lens", 0.3, 0.06, 1.2, 0.42,
             font_size=20, bold=True, color=WHITE, font_name="Calibri Light")
    add_text(slide, "&", 1.25, 0.06, 0.4, 0.42,
             font_size=20, bold=True, color=GOLD)
    add_text(slide, "Light", 1.52, 0.06, 1.2, 0.42,
             font_size=20, bold=True, color=WHITE, font_name="Calibri Light")


# ─── SLIDE BUILDERS ───────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Gold top stripe
    add_rect(slide, 0, 0, 13.33, 0.12, fill_color=GOLD)
    # Gold bottom stripe
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=GOLD)

    # Left gold vertical accent
    add_rect(slide, 0, 0.12, 0.12, 7.26, fill_color=BG_CARD)

    # Background card
    add_rect(slide, 1.2, 1.0, 11.0, 5.5, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1.5))

    # Camera icon placeholder text
    add_text(slide, "📸", 6.0, 1.2, 1.5, 1.0, font_size=40, color=GOLD,
             align=PP_ALIGN.CENTER)

    # Title
    add_text(slide, "LENS & LIGHT STUDIO",
             1.5, 2.1, 10.5, 0.85,
             font_size=44, bold=True, color=GOLD_LIGHT,
             align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # Subtitle
    add_text(slide, "Full-Stack Photography Business Web Application",
             1.5, 2.85, 10.5, 0.5,
             font_size=20, bold=False, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Divider
    add_rect(slide, 4.5, 3.4, 4.33, 0.04, fill_color=GOLD)

    # Info lines
    add_text(slide, "Submitted by: Rohit Rajput",
             1.5, 3.55, 10.5, 0.35, font_size=14, color=GREY,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Academic Year: 2025–2026",
             1.5, 3.88, 10.5, 0.35, font_size=14, color=GREY,
             align=PP_ALIGN.CENTER)
    add_text(slide, "Dept. of Computer Science & Engineering",
             1.5, 4.18, 10.5, 0.35, font_size=14, color=GREY,
             align=PP_ALIGN.CENTER)

    # Tech tags
    for i, tag in enumerate(["Python", "Flask", "SQLite", "HTML5", "CSS3", "JavaScript"]):
        x = 1.9 + i * 1.6
        add_rect(slide, x, 5.0, 1.35, 0.38, fill_color=BG_DARK,
                 line_color=GOLD, line_width=Pt(0.8))
        add_text(slide, tag, x+0.05, 5.04, 1.25, 0.32,
                 font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Slide number
    add_text(slide, "01", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_02_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Agenda", "What we will cover today")

    items = [
        ("01", "Introduction & Problem Statement"),
        ("02", "Objectives of the Project"),
        ("03", "Technology Stack"),
        ("04", "System Architecture"),
        ("05", "Database Design"),
        ("06", "Key Features — Public Website"),
        ("07", "Key Features — Admin Panel"),
        ("08", "Testing & Results"),
        ("09", "Future Scope"),
        ("10", "Conclusion"),
    ]

    cols = [items[:5], items[5:]]
    for ci, col in enumerate(cols):
        for ri, (num, text) in enumerate(col):
            x = 0.4 + ci * 6.5
            y = 1.65 + ri * 0.98
            add_rect(slide, x, y, 6.2, 0.78, fill_color=BG_CARD,
                     line_color=GOLD, line_width=Pt(0.8))
            add_text(slide, num, x+0.1, y+0.15, 0.6, 0.5,
                     font_size=18, bold=True, color=GOLD)
            add_text(slide, text, x+0.7, y+0.18, 5.3, 0.45,
                     font_size=14, bold=False, color=WHITE)

    add_text(slide, "10", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_03_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Introduction & Problem Statement",
                 "Why this project was needed")

    add_text(slide, "The Challenge",
             0.4, 1.65, 6.0, 0.4, font_size=17, bold=True, color=GOLD)

    problems = [
        "No centralised professional online presence",
        "Booking managed via WhatsApp / email — error prone",
        "Prices & services cannot be updated without a developer",
        "Client enquiries scattered across multiple platforms",
        "No business analytics or booking status tracking",
        "Client data stored on third-party servers (privacy risk)",
    ]
    for i, p in enumerate(problems):
        y = 2.05 + i * 0.72
        add_rect(slide, 0.4, y, 6.0, 0.58, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(0.7))
        add_text(slide, "✗  " + p, 0.55, y+0.1, 5.7, 0.4,
                 font_size=12.5, color=WHITE)

    # Right side — Our Solution
    add_text(slide, "Our Solution",
             7.0, 1.65, 6.0, 0.4, font_size=17, bold=True, color=GOLD)

    add_rect(slide, 7.0, 2.05, 6.0, 5.25, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))

    add_text(slide,
             "Lens & Light Studio is a complete,\n"
             "self-hosted, full-stack web application\n"
             "that gives photographers:\n\n"
             "✔  A professional portfolio website\n"
             "✔  An online booking system\n"
             "✔  A contact & inquiry manager\n"
             "✔  A secure admin dashboard\n"
             "✔  Full content management (CRUD)\n"
             "✔  Zero recurring SaaS fees",
             7.2, 2.2, 5.6, 4.9, font_size=14, color=WHITE)

    add_text(slide, "03", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_04_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Objectives of the Project", "")

    objs = [
        ("🎯", "Full-Stack Web App", "Using Flask, SQLAlchemy, SQLite, HTML5, CSS3, JS"),
        ("🖼️", "Dynamic Gallery", "Filterable masonry grid with lightbox viewer"),
        ("📅", "Online Booking", "Multi-field booking form saved to database"),
        ("🔐", "Secure Admin Panel", "bcrypt authentication & session management"),
        ("📩", "Contact Management", "Inquiry capture, read tracking & management"),
        ("📱", "Responsive Design", "Works on all screens from 320px to 1920px"),
        ("⚙️", "Full CRUD", "Gallery, Services, Bookings, Testimonials, Inquiries"),
        ("🚀", "Extensible", "Architecture ready for payments, email, calendar"),
    ]

    for i, (icon, title, desc) in enumerate(objs):
        row = i // 4
        col = i % 4
        x = 0.25 + col * 3.22
        y = 1.55 + row * 2.65
        add_rect(slide, x, y, 3.0, 2.35, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(1))
        add_text(slide, icon, x+1.1, y+0.12, 0.8, 0.55, font_size=24)
        add_text(slide, title, x+0.1, y+0.65, 2.8, 0.42,
                 font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x+0.1, y+1.05, 2.8, 1.1,
                 font_size=11, color=GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "04", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_05_techstack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Technology Stack", "Tools and frameworks used")

    layers = [
        ("FRONTEND", [
            ("HTML5", "Semantic structure, SEO-friendly markup"),
            ("CSS3",  "Custom design system, glassmorphism, animations"),
            ("JavaScript ES6+", "IntersectionObserver, Fetch API, Lightbox"),
            ("Google Fonts", "Playfair Display (headings) + Inter (body)"),
        ]),
        ("BACKEND", [
            ("Python 3.13", "Core programming language"),
            ("Flask 3.x", "Micro web framework — routing, templating"),
            ("Jinja2", "Server-side HTML templating engine"),
            ("Werkzeug", "File upload security & filename sanitisation"),
        ]),
        ("DATABASE & SECURITY", [
            ("SQLite", "Zero-config embedded relational database"),
            ("SQLAlchemy ORM", "Pythonic database abstraction layer"),
            ("Flask-Login", "Session management & route protection"),
            ("bcrypt", "Adaptive password hashing algorithm"),
        ]),
    ]

    for ci, (layer_title, items) in enumerate(layers):
        x = 0.3 + ci * 4.35
        add_rect(slide, x, 1.5, 4.1, 5.65, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(1.2))
        add_rect(slide, x, 1.5, 4.1, 0.45, fill_color=GOLD)
        add_text(slide, layer_title, x+0.15, 1.55, 3.8, 0.38,
                 font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        for ri, (name, desc) in enumerate(items):
            y = 2.05 + ri * 1.2
            add_rect(slide, x+0.18, y, 3.72, 1.08, fill_color=BG_DARK,
                     line_color=hex_rgb("#2a2a3a"), line_width=Pt(0.5))
            add_text(slide, name, x+0.3, y+0.08, 3.5, 0.4,
                     font_size=13, bold=True, color=GOLD_LIGHT)
            add_text(slide, desc, x+0.3, y+0.48, 3.5, 0.52,
                     font_size=11, color=GREY)

    add_text(slide, "05", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_06_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "System Architecture", "Three-Tier MVC Pattern")

    # Arrow function
    def arrow(y_center):
        add_rect(slide, 6.46, y_center-0.06, 0.4, 0.12, fill_color=GOLD)
        add_text(slide, "▶", 6.75, y_center-0.17, 0.3, 0.32,
                 font_size=13, color=GOLD, align=PP_ALIGN.CENTER)

    tiers = [
        ("PRESENTATION LAYER  (View)", BG_CARD, GOLD, [
            "HTML5 Templates (Jinja2)",
            "CSS3 Design System",
            "Vanilla JavaScript",
            "Google Fonts",
        ]),
        ("APPLICATION LAYER  (Controller)", BG_CARD, GOLD_LIGHT, [
            "Flask Blueprints (public + admin)",
            "Route Handlers / Business Logic",
            "Flask-Login Authentication",
            "File Upload Handler",
        ]),
        ("DATA LAYER  (Model)", BG_CARD, GOLD, [
            "SQLAlchemy ORM Models",
            "SQLite Database (6 tables)",
            "bcrypt Password Hashing",
            "Data Validation",
        ]),
    ]

    for i, (title, bg, col, items) in enumerate(tiers):
        y = 1.55 + i * 1.9
        add_rect(slide, 0.3, y, 6.0, 1.65, fill_color=bg,
                 line_color=col, line_width=Pt(1.2))
        add_text(slide, title, 0.5, y+0.1, 5.6, 0.42,
                 font_size=13, bold=True, color=col)
        for j, item in enumerate(items):
            add_text(slide, "  •  " + item, 0.5, y+0.5 + j*0.27,
                     5.6, 0.28, font_size=11, color=WHITE)
        if i < 2:
            arrow(y + 1.65)

    # Right panel — Request/Response Flow
    add_rect(slide, 7.0, 1.55, 6.0, 5.7, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))
    add_text(slide, "Request / Response Flow",
             7.15, 1.65, 5.7, 0.42, font_size=14, bold=True, color=GOLD)

    steps = [
        ("1", "Client sends HTTP request (e.g. GET /gallery)"),
        ("2", "Flask Router matches URL to route handler"),
        ("3", "Route handler queries database via SQLAlchemy"),
        ("4", "SQLAlchemy executes SQL against SQLite"),
        ("5", "Data returned to route handler"),
        ("6", "Jinja2 renders HTML template with data"),
        ("7", "Flask sends HTML response to browser"),
        ("8", "Browser renders HTML, applies CSS, runs JS"),
    ]
    for i, (n, s) in enumerate(steps):
        y = 2.15 + i * 0.62
        add_rect(slide, 7.15, y, 0.38, 0.38, fill_color=GOLD)
        add_text(slide, n, 7.15, y+0.03, 0.38, 0.33,
                 font_size=12, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, s, 7.62, y+0.05, 5.2, 0.35,
                 font_size=11, color=WHITE)

    add_text(slide, "06", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_07_database(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Database Design", "6 tables — SQLite via SQLAlchemy ORM")

    tables = [
        ("USERS", ["id (PK)", "username (UNIQUE)", "email (UNIQUE)", "password_hash", "is_admin"]),
        ("GALLERY", ["id (PK)", "title", "category", "filename", "description", "featured", "created_at"]),
        ("SERVICES", ["id (PK)", "name", "price", "description", "features_json", "popular", "duration"]),
        ("BOOKINGS", ["id (PK)", "name", "email", "phone", "service_id (FK)", "event_date", "status"]),
        ("TESTIMONIALS", ["id (PK)", "client_name", "rating", "review", "service", "avatar_initials"]),
        ("INQUIRIES", ["id (PK)", "name", "email", "subject", "message", "read", "created_at"]),
    ]

    positions = [
        (0.3, 1.55), (4.55, 1.55), (8.8, 1.55),
        (0.3, 4.55), (4.55, 4.55), (8.8, 4.55),
    ]

    for (x, y), (name, cols) in zip(positions, tables):
        w, h = 4.1, 2.8
        add_rect(slide, x, y, w, h, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(1))
        add_rect(slide, x, y, w, 0.42, fill_color=GOLD)
        add_text(slide, name, x+0.1, y+0.06, w-0.2, 0.32,
                 font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        for ri, col in enumerate(cols):
            bg = BG_DARK if ri % 2 == 0 else hex_rgb("#14141F")
            add_rect(slide, x+0.05, y+0.46+ri*0.36, w-0.1, 0.34, fill_color=bg)
            add_text(slide, col, x+0.18, y+0.5+ri*0.36, w-0.25, 0.28,
                     font_size=10.5, color=WHITE if "PK" not in col and "FK" not in col else GOLD_LIGHT)

    # Relationship note
    add_text(slide, "BOOKINGS.service_id  →  SERVICES.id  (Many-to-One)",
             0.3, 7.1, 9.0, 0.3, font_size=11, color=GREY)

    add_text(slide, "07", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_08_public_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Key Features — Public Website",
                 "7 pages for prospective clients")

    pages = [
        ("🏠  Home",       ["Hero section with animated stats", "Featured gallery (6 photos)", "Service preview cards", "Client testimonials"]),
        ("🖼️  Gallery",    ["12 real portfolio photos", "Category filter buttons", "Masonry grid layout", "Lightbox full-screen view"]),
        ("💼  Services",   ["6 pricing packages", "Most Popular badge", "Features list per package", "Book Now CTA buttons"]),
        ("📅  Booking",    ["Multi-field booking form", "Visual package selector", "Event date & type fields", "Server-side validation"]),
        ("📞  Contact",    ["Contact info cards", "Message form with subject", "Success/error feedback", "Saved to DB as inquiry"]),
        ("⭐  Reviews",    ["Star-rated testimonials", "Client avatar initials", "Service label per review", "Admin-managed content"]),
        ("👤  About",      ["Photographer biography", "Professional stats", "Skills & achievements", "Brand story section"]),
    ]

    for i, (title, bullets) in enumerate(pages):
        row = i // 4
        col = i % 4
        if i == 4:
            col = 0
        if i >= 4:
            col = i - 4
        x = 0.2 + col * 3.28
        y = 1.6 + row * 2.85
        if i >= 4:
            y = 4.45
        w, h = 3.1, 2.6
        add_rect(slide, x, y, w, h, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(0.9))
        add_text(slide, title, x+0.12, y+0.1, w-0.2, 0.42,
                 font_size=13, bold=True, color=GOLD)
        for bi, b in enumerate(bullets):
            add_text(slide, "  ▸  " + b, x+0.1, y+0.55+bi*0.48,
                     w-0.2, 0.42, font_size=11, color=WHITE)

    add_text(slide, "08", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_09_admin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Key Features — Admin Panel",
                 "Secure management dashboard for the photographer")

    sections = [
        ("🔐  Login", "bcrypt-secured login\nFlask-Login sessions\nAuto-redirect if logged in\nLogout terminates session"),
        ("📊  Dashboard", "8 live stat cards\nTotal / Pending / Confirmed\nUnread inquiry count\nRecent activity tables"),
        ("🖼️  Gallery Mgr", "Upload with category tag\nFeatured toggle for homepage\nDelete with file cleanup\nPhoto preview grid"),
        ("📅  Bookings Mgr", "View all bookings\nStatus: Pending / Confirmed / Cancelled\nDelete records\nFilter by status"),
        ("💼  Services Mgr", "Add / Edit / Delete packages\nPrice, duration, features\nMost Popular badge toggle\nReal-time on public page"),
        ("📩  Inquiries Mgr", "Read/Unread highlighting\nMark Read (AJAX, no reload)\nDelete messages\nTimestamp tracking"),
    ]

    for i, (title, desc) in enumerate(sections):
        row = i // 3
        col = i % 3
        x = 0.3 + col * 4.35
        y = 1.55 + row * 2.85
        add_rect(slide, x, y, 4.1, 2.6, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(1))
        add_rect(slide, x, y, 4.1, 0.44, fill_color=hex_rgb("#1a1a2e"))
        add_text(slide, title, x+0.15, y+0.06, 3.8, 0.34,
                 font_size=13, bold=True, color=GOLD)
        lines = desc.split('\n')
        for li, line in enumerate(lines):
            add_text(slide, "  •  " + line, x+0.15, y+0.55+li*0.49,
                     3.8, 0.44, font_size=12, color=WHITE)

    add_text(slide, "09", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_10_ui_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "UI Design System", "Dark Gold Glassmorphism Aesthetic")

    # Colour palette
    add_text(slide, "Colour Palette", 0.4, 1.6, 4.0, 0.38,
             font_size=15, bold=True, color=GOLD)

    palette = [
        ("#0a0a0f", "Background"),
        ("#111118", "Card BG"),
        ("#c9a84c", "Gold Accent"),
        ("#e8c97a", "Gold Light"),
        ("#f0ede8", "Off-White"),
        ("#9b9592", "Muted Grey"),
    ]
    for i, (hx, name) in enumerate(palette):
        y = 2.0 + i * 0.72
        add_rect(slide, 0.4, y, 0.65, 0.55, fill_color=hex_rgb(hx),
                 line_color=GOLD, line_width=Pt(0.5))
        add_text(slide, hx, 1.15, y+0.03, 1.5, 0.28, font_size=11, color=GREY)
        add_text(slide, name, 1.15, y+0.28, 1.5, 0.25, font_size=12,
                 bold=True, color=WHITE)

    # Typography
    add_text(slide, "Typography", 3.3, 1.6, 4.0, 0.38,
             font_size=15, bold=True, color=GOLD)
    add_text(slide, "Playfair Display", 3.3, 2.05, 4.0, 0.55,
             font_size=26, bold=True, color=WHITE, font_name="Georgia")
    add_text(slide, "Headings, Logo — Elegant Serif", 3.3, 2.55, 4.0, 0.32,
             font_size=12, color=GREY)
    add_text(slide, "Inter — Body & UI Text", 3.3, 3.0, 4.0, 0.45,
             font_size=18, bold=False, color=WHITE)
    add_text(slide, "Body text, labels, buttons — Clean Sans-Serif", 3.3, 3.42, 4.0, 0.32,
             font_size=12, color=GREY)

    # Design features
    add_text(slide, "Design Features", 3.3, 4.0, 4.0, 0.38,
             font_size=15, bold=True, color=GOLD)
    features = [
        "Glassmorphism — backdrop blur + semi-transparent cards",
        "CSS Custom Properties for consistent design tokens",
        "Scroll Reveal with IntersectionObserver API",
        "Counter animation on homepage statistics",
        "Hover lift effect on all interactive cards",
        "Hamburger menu for mobile navigation",
    ]
    for i, f in enumerate(features):
        add_text(slide, "  ▸  " + f, 3.3, 4.42+i*0.44, 3.85, 0.4,
                 font_size=11.5, color=WHITE)

    # Breakpoints
    add_text(slide, "Responsive Breakpoints", 7.5, 1.6, 5.6, 0.38,
             font_size=15, bold=True, color=GOLD)
    bps = [
        ("Mobile", "< 768px", "Single column, hamburger nav"),
        ("Tablet", "768–1024px", "Two columns, medium layout"),
        ("Desktop", "> 1024px", "Full layout, three-col gallery"),
    ]
    for i, (device, bp, note) in enumerate(bps):
        y = 2.05 + i * 1.05
        add_rect(slide, 7.5, y, 5.6, 0.9, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(0.8))
        add_text(slide, device, 7.65, y+0.06, 1.5, 0.35,
                 font_size=13, bold=True, color=GOLD)
        add_text(slide, bp, 9.3, y+0.06, 1.5, 0.35,
                 font_size=12, bold=True, color=GOLD_LIGHT)
        add_text(slide, note, 7.65, y+0.44, 5.2, 0.35,
                 font_size=11, color=GREY)

    # Animations
    add_text(slide, "JavaScript Interactions", 7.5, 5.25, 5.6, 0.38,
             font_size=15, bold=True, color=GOLD)
    anis = [
        "Lightbox gallery viewer",
        "Async Mark-as-Read (Fetch API, no reload)",
        "Sticky nav glass effect on scroll",
        "Package card selector for booking form",
    ]
    for i, a in enumerate(anis):
        add_text(slide, "  ✦  " + a, 7.5, 5.65+i*0.44, 5.6, 0.4,
                 font_size=12, color=WHITE)

    add_text(slide, "10", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_11_testing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Testing & Results",
                 "7 levels of testing — all passed")

    # Summary stats
    stats = [
        ("57", "Total Test Cases"),
        ("100%", "Pass Rate"),
        ("4", "Browsers Tested"),
        ("8", "Screen Sizes"),
    ]
    for i, (val, label) in enumerate(stats):
        x = 0.3 + i * 3.22
        add_rect(slide, x, 1.55, 3.0, 1.15, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(1))
        add_text(slide, val, x, 1.62, 3.0, 0.6,
                 font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, label, x, 2.2, 3.0, 0.38,
                 font_size=12, color=GREY, align=PP_ALIGN.CENTER)

    # Test types table
    headers = ["Testing Level", "Cases", "Result"]
    rows = [
        ["Unit Testing — Public Routes", "12", "✔ PASS"],
        ["Unit Testing — Admin Routes", "19", "✔ PASS"],
        ["Integration Testing", "10", "✔ PASS"],
        ["System Testing (E2E)", "10", "✔ PASS"],
        ["User Acceptance Testing", "6 criteria", "✔ PASS"],
        ["Browser Compatibility (4 browsers)", "10 features", "✔ PASS"],
        ["Responsive Design (8 sizes)", "All breakpoints", "✔ PASS"],
        ["Security Testing", "8 tests", "✔ PASS"],
    ]

    col_widths = [5.8, 1.8, 1.8]
    col_x = [0.3, 6.1, 7.9]
    header_y = 2.85

    # Header row
    for ci, (hdr, w, x) in enumerate(zip(headers, col_widths, col_x)):
        add_rect(slide, x, header_y, w, 0.42, fill_color=GOLD)
        add_text(slide, hdr, x+0.1, header_y+0.07, w-0.15, 0.3,
                 font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(rows):
        y = 3.3 + ri * 0.47
        bg = BG_CARD if ri % 2 == 0 else BG_DARK
        for ci, (cell, w, x) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, x, y, w, 0.43, fill_color=bg,
                     line_color=hex_rgb("#2a2a3a"), line_width=Pt(0.3))
            clr = hex_rgb("#4CAF50") if "✔" in cell else WHITE
            add_text(slide, cell, x+0.1, y+0.09, w-0.15, 0.28,
                     font_size=12, color=clr,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Security highlights
    add_rect(slide, 9.9, 2.85, 3.1, 4.45, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))
    add_text(slide, "🔐  Security", 10.0, 2.95, 2.9, 0.38,
             font_size=13, bold=True, color=GOLD)
    sec = [
        "bcrypt password hashing",
        "@login_required on all admin routes",
        "File extension whitelist",
        "Secure filename sanitisation",
        "16 MB upload size limit",
        "Parameterised SQL queries",
        "HMAC-signed session cookies",
    ]
    for i, s in enumerate(sec):
        add_text(slide, "  ✔  " + s, 10.0, 3.4+i*0.48, 2.9, 0.42,
                 font_size=11, color=WHITE)

    add_text(slide, "11", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_12_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Future Scope",
                 "Planned enhancements and extensions")

    future = [
        ("💳", "Payment Integration",    "Stripe / Razorpay deposit at booking time"),
        ("📧", "Email Notifications",    "Auto-confirm emails via Flask-Mail"),
        ("📆", "Availability Calendar",  "Prevent double-bookings, Google Calendar sync"),
        ("👤", "Client Portal",          "Login to view booking status & photo gallery"),
        ("📝", "Blog / SEO Content",     "Improve organic search rankings"),
        ("📊", "Advanced Analytics",     "Revenue tracking, seasonal trends"),
        ("📱", "Mobile App",             "React Native companion for on-the-go management"),
        ("☁️", "Cloud Storage",          "AWS S3 / Cloudinary for unlimited photo storage"),
        ("🌐", "PWA Support",            "Offline gallery, Add-to-Home-Screen"),
        ("👥", "Multi-Photographer",     "Studio with multiple photographer profiles"),
    ]

    for i, (icon, title, desc) in enumerate(future):
        row = i // 5
        col = i % 5
        x = 0.2 + col * 2.6
        y = 1.55 + row * 2.75
        add_rect(slide, x, y, 2.45, 2.5, fill_color=BG_CARD,
                 line_color=GOLD, line_width=Pt(0.8))
        add_text(slide, icon, x+0.85, y+0.12, 0.8, 0.55, font_size=24,
                 align=PP_ALIGN.CENTER)
        add_text(slide, title, x+0.1, y+0.68, 2.25, 0.55,
                 font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x+0.1, y+1.22, 2.25, 1.0,
                 font_size=10.5, color=GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "12", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_13_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    logo_text(slide)
    slide_header(slide, "Conclusion", "What was achieved")

    # Left — achievements
    add_rect(slide, 0.3, 1.55, 6.1, 5.65, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))
    add_text(slide, "All 10 Objectives Achieved",
             0.45, 1.65, 5.8, 0.42, font_size=15, bold=True, color=GOLD)

    achieved = [
        "Full-stack Flask + SQLAlchemy + SQLite app",
        "7 public pages with premium dark-gold UI",
        "12 gallery photos across 6 categories",
        "Filterable masonry gallery + lightbox",
        "Online booking system — saved to database",
        "Contact inquiry management system",
        "Secure bcrypt admin panel",
        "Full CRUD for all 5 content types",
        "Fully responsive — 320px to 1920px",
        "57 test cases — 100% pass rate",
    ]
    for i, a in enumerate(achieved):
        y = 2.12 + i * 0.5
        add_text(slide, "  ✔  " + a, 0.5, y, 5.7, 0.44,
                 font_size=12.5, color=WHITE)

    # Right — key learnings
    add_rect(slide, 6.7, 1.55, 6.3, 2.6, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))
    add_text(slide, "Key Learnings", 6.85, 1.65, 5.9, 0.4,
             font_size=15, bold=True, color=GOLD)
    learnings = [
        "Flask blueprint architecture & App Factory pattern",
        "SQLAlchemy ORM — models, relationships, queries",
        "bcrypt security & session-based authentication",
        "CSS glassmorphism, Custom Properties, animations",
        "JavaScript Fetch API & IntersectionObserver",
        "Full software engineering lifecycle (SDLC)",
    ]
    for i, l in enumerate(learnings):
        add_text(slide, "  ▸  " + l, 6.85, 2.12+i*0.34, 6.0, 0.3,
                 font_size=11.5, color=WHITE)

    # Cost comparison
    add_rect(slide, 6.7, 4.3, 6.3, 2.9, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1))
    add_text(slide, "Cost Comparison", 6.85, 4.42, 5.9, 0.4,
             font_size=15, bold=True, color=GOLD)
    costs = [
        ("Pixieset (SaaS)",      "₹2,000–₹8,500/month"),
        ("Squarespace",          "₹1,500–₹3,400/month"),
        ("Custom Dev (agency)",  "₹50,000–₹2,00,000 one-time"),
        ("This Project",         "₹400–₹800/month (hosting only)"),
    ]
    for i, (name, cost) in enumerate(costs):
        y = 4.9 + i * 0.52
        clr = GOLD if i == 3 else GREY
        bold = True if i == 3 else False
        add_text(slide, name, 6.85, y, 3.5, 0.4, font_size=12,
                 bold=bold, color=clr)
        add_text(slide, cost, 10.3, y, 2.5, 0.4, font_size=12,
                 bold=bold, color=clr, align=PP_ALIGN.RIGHT)

    add_text(slide, "13", 12.9, 7.1, 0.4, 0.3,
             font_size=9, color=GREY, align=PP_ALIGN.RIGHT)


def slide_14_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_rect(slide, 0, 0, 13.33, 0.12, fill_color=GOLD)
    add_rect(slide, 0, 7.38, 13.33, 0.12, fill_color=GOLD)

    add_rect(slide, 2.0, 1.0, 9.33, 5.6, fill_color=BG_CARD,
             line_color=GOLD, line_width=Pt(1.5))

    add_text(slide, "📸", 6.15, 1.2, 1.5, 1.0, font_size=48,
             color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "Thank You",
             2.0, 2.1, 9.33, 0.85, font_size=52, bold=True,
             color=GOLD_LIGHT, align=PP_ALIGN.CENTER, font_name="Calibri Light")

    add_rect(slide, 5.0, 3.05, 3.33, 0.05, fill_color=GOLD)

    add_text(slide, "Lens & Light Studio",
             2.0, 3.18, 9.33, 0.5, font_size=20, bold=False,
             color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Full-Stack Photography Business Web Application",
             2.0, 3.65, 9.33, 0.4, font_size=14, bold=False,
             color=GREY, align=PP_ALIGN.CENTER)

    add_text(slide, "Submitted by: Rohit Rajput",
             2.0, 4.2, 9.33, 0.4, font_size=14,
             color=GREY, align=PP_ALIGN.CENTER)

    # Tech tags
    for i, tag in enumerate(["Python", "Flask", "SQLAlchemy", "SQLite", "HTML5", "CSS3", "JS"]):
        total = 7
        total_width = total * 1.35 + (total-1) * 0.15
        start_x = (13.33 - total_width) / 2
        x = start_x + i * 1.5
        add_rect(slide, x, 5.0, 1.35, 0.42, fill_color=BG_DARK,
                 line_color=GOLD, line_width=Pt(0.8))
        add_text(slide, tag, x+0.05, 5.05, 1.25, 0.32,
                 font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    add_text(slide, "Questions & Discussion Welcome",
             2.0, 5.7, 9.33, 0.4, font_size=14,
             color=GREY, align=PP_ALIGN.CENTER)


# ─── MAIN ─────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide_01_title(prs)         ; print("  [1/14] Title")
    slide_02_agenda(prs)        ; print("  [2/14] Agenda")
    slide_03_intro(prs)         ; print("  [3/14] Introduction")
    slide_04_objectives(prs)    ; print("  [4/14] Objectives")
    slide_05_techstack(prs)     ; print("  [5/14] Tech Stack")
    slide_06_architecture(prs)  ; print("  [6/14] Architecture")
    slide_07_database(prs)      ; print("  [7/14] Database")
    slide_08_public_features(prs); print(" [8/14] Public Features")
    slide_09_admin(prs)         ; print("  [9/14] Admin Panel")
    slide_10_ui_design(prs)     ; print(" [10/14] UI Design")
    slide_11_testing(prs)       ; print(" [11/14] Testing")
    slide_12_future(prs)        ; print(" [12/14] Future Scope")
    slide_13_conclusion(prs)    ; print(" [13/14] Conclusion")
    slide_14_thankyou(prs)      ; print(" [14/14] Thank You")

    out = r"c:\Users\Rohit Rajput\OneDrive\Documents\New project\Lens_Light_Studio_Presentation.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")

if __name__ == "__main__":
    build()
